#!/usr/bin/env python3
"""
HTML Crop Tables → PowerPoint Converter
========================================
Parses HTML crop growth stage tables and generates PPTX presentations.

Usage:
    python html_to_pptx.py --html-dir ./crops --images-dir ./assets/images/crops --output-dir ./pptx_output

Each HTML file produces one PPTX file with a single landscape slide containing:
- Crop title
- 10 stage images in a row
- BBCH codes under each image
- Stage descriptions
- Empty "Your Product" row for client to fill in

Or use --single-file mode to combine all crops into one PPTX (one slide per crop).
"""

import os
import re
import glob
import argparse
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ─── DESIGN CONSTANTS ──────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)

# Colors (no # prefix)
CLR_TITLE = RGBColor(0x6C, 0x74, 0x66)       # muted green-gray
CLR_SUBTITLE = RGBColor(0x9D, 0xA3, 0x9A)    # lighter gray
CLR_LABEL = RGBColor(0x6C, 0x74, 0x66)        # same as title
CLR_BBCH = RGBColor(0x4A, 0x4F, 0x45)         # dark text
CLR_DESC = RGBColor(0x5A, 0x5E, 0x54)         # medium text
CLR_PLACEHOLDER = RGBColor(0xC8, 0xCE, 0xC0)  # light gray placeholder
CLR_BORDER = RGBColor(0xC8, 0xCE, 0xC0)       # table border
CLR_HEADER_BG = RGBColor(0xF2, 0xF4, 0xF0)    # subtle header background
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_FOOTER = RGBColor(0xB0, 0xB5, 0xAA)

# Layout
MARGIN_LEFT = Inches(0.4)
MARGIN_RIGHT = Inches(0.4)
MARGIN_TOP = Inches(0.3)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

NUM_STAGES = 10
LABEL_COL_WIDTH = Inches(1.3)
STAGE_AREA_WIDTH = CONTENT_WIDTH - LABEL_COL_WIDTH


def parse_html(html_path):
    """Parse a crop HTML file and extract structured data."""
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    data = {}

    # Title from h1
    h1 = soup.find('h1')
    data['title'] = h1.get_text(strip=True) if h1 else Path(html_path).stem

    # Extract crop slug from image paths
    image_row = soup.find('tr', class_='image-row')
    if image_row:
        imgs = image_row.find_all('img')
        data['images'] = []
        data['image_alts'] = []
        for img in imgs:
            src = img.get('src', '')
            alt = img.get('alt', '')
            data['images'].append(src)
            data['image_alts'].append(alt)

    # Extract BBCH codes
    data_rows = soup.find_all('tr', class_='data-row')
    data['bbch_codes'] = []
    data['descriptions'] = []

    for row in data_rows:
        cells = row.find_all('td')
        if not cells:
            continue

        label_cell = cells[0]
        label_text = label_cell.get_text(strip=True).lower()

        values = [cell.get_text(strip=True) for cell in cells[1:]]

        if 'bbch' in label_text or 'stage' == label_text.split()[-1] if label_text.split() else False:
            # Check if cells have bbch class
            bbch_cells = row.find_all('td', class_='bbch')
            if bbch_cells:
                data['bbch_codes'] = [c.get_text(strip=True) for c in bbch_cells]
            elif values:
                data['bbch_codes'] = values
        elif 'description' in label_text or 'описан' in label_text:
            data['descriptions'] = values
        elif 'product' in label_text or 'препарат' in label_text:
            pass  # Skip product row, we'll create empty cells

    # Fallback: if bbch not found via label, try finding by class
    if not data['bbch_codes']:
        for row in data_rows:
            bbch_cells = row.find_all('td', class_='bbch')
            if bbch_cells:
                data['bbch_codes'] = [c.get_text(strip=True) for c in bbch_cells]
                break

    # Extract crop slug from image paths to find local images
    if data.get('images'):
        # e.g. "../assets/images/crops/alfalfa/alfalfa_stage_1.png"
        first_img = data['images'][0]
        # Try to extract the crop folder name
        parts = first_img.replace('\\', '/').split('/')
        for i, p in enumerate(parts):
            if p == 'crops' and i + 1 < len(parts):
                data['crop_slug'] = parts[i + 1]
                break
        if 'crop_slug' not in data:
            # Fallback: derive from filename
            basename = Path(first_img).stem
            # e.g. "alfalfa_stage_1" → "alfalfa"
            match = re.match(r'(.+?)_stage_\d+', basename)
            if match:
                data['crop_slug'] = match.group(1)
            else:
                data['crop_slug'] = Path(html_path).stem

    return data


def resolve_image_path(crop_slug, stage_num, images_dir):
    """Find the actual image file for a given crop and stage."""
    # Try standard naming: {crop_slug}_stage_{n}.png
    candidates = [
        os.path.join(images_dir, crop_slug, f'{crop_slug}_stage_{stage_num}.png'),
        os.path.join(images_dir, crop_slug, f'{crop_slug}_stage_{stage_num}.jpg'),
    ]
    for c in candidates:
        if os.path.isfile(c):
            return c
    return None


def add_crop_slide(prs, crop_data, images_dir, include_footer=True):
    """Add a single crop slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CLR_WHITE

    crop_slug = crop_data.get('crop_slug', '')

    # Determine which stages actually have images
    existing_stages = []
    for i in range(1, 11):
        img_path = resolve_image_path(crop_slug, i, images_dir)
        if img_path:
            existing_stages.append(i)

    # Fallback: if no images found, use all 10
    if not existing_stages:
        existing_stages = list(range(1, 11))

    num_stages = len(existing_stages)
    stage_col_w = STAGE_AREA_WIDTH / num_stages

    y_cursor = MARGIN_TOP

    # ─── TITLE ──────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, y_cursor,
        CONTENT_WIDTH * 0.75, Inches(0.45)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = crop_data['title']
    p.font.size = Pt(22)
    p.font.color.rgb = CLR_TITLE
    p.font.name = 'Segoe UI'
    p.font.bold = False

    # Subtitle on the right
    sub_box = slide.shapes.add_textbox(
        MARGIN_LEFT + CONTENT_WIDTH * 0.75, y_cursor,
        CONTENT_WIDTH * 0.25, Inches(0.45)
    )
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = 'Botanical Growth Stages'
    p2.font.size = Pt(10)
    p2.font.color.rgb = CLR_SUBTITLE
    p2.font.name = 'Segoe UI'
    p2.alignment = PP_ALIGN.RIGHT

    y_cursor += Inches(0.55)

    # ─── IMAGES ROW ─────────────────────────────────────────
    img_row_height = Inches(2.4)
    img_x_start = MARGIN_LEFT + LABEL_COL_WIDTH

    images_found = 0

    for col_idx, stage_num in enumerate(existing_stages):
        img_path = resolve_image_path(crop_slug, stage_num, images_dir)

        if img_path:
            images_found += 1
            max_img_w = stage_col_w - Inches(0.05)
            max_img_h = img_row_height - Inches(0.2)

            try:
                pic = slide.shapes.add_picture(
                    img_path,
                    int(img_x_start + stage_col_w * col_idx + Inches(0.025)),
                    int(y_cursor),
                    int(max_img_w),
                    int(max_img_h)
                )
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pil_img:
                    orig_w, orig_h = pil_img.size

                aspect = orig_w / orig_h
                target_w = max_img_w
                target_h = int(target_w / aspect)
                if target_h > max_img_h:
                    target_h = max_img_h
                    target_w = int(target_h * aspect)

                pic.width = int(target_w)
                pic.height = int(target_h)

                col_center_x = img_x_start + stage_col_w * col_idx + stage_col_w / 2
                pic.left = int(col_center_x - target_w / 2)
                pic.top = int(y_cursor + img_row_height - target_h)

            except Exception as e:
                print(f"  Warning: Could not add image {img_path}: {e}")

    y_cursor += img_row_height + Inches(0.05)

    # ─── DATA TABLE ─────────────────────────────────────────
    num_cols = num_stages + 1  # label column + stage columns
    num_rows = 3

    table_height = Inches(1.8)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        int(MARGIN_LEFT), int(y_cursor),
        int(CONTENT_WIDTH), int(table_height)
    )
    table = table_shape.table

    # Column widths
    table.columns[0].width = int(LABEL_COL_WIDTH)
    for i in range(1, num_cols):
        table.columns[i].width = int(stage_col_w)

    # Style helper
    def style_cell(cell, text, font_size=8, bold=False, color=CLR_BBCH,
                   align=PP_ALIGN.CENTER, fill_color=None):
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.alignment = align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Margins
        cell.text_frame.margin_left = Pt(3)
        cell.text_frame.margin_right = Pt(3)
        cell.text_frame.margin_top = Pt(3)
        cell.text_frame.margin_bottom = Pt(3)
        cell.text_frame.word_wrap = True

        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
        else:
            cell.fill.background()

    # Row 0: BBCH Stage
    style_cell(table.cell(0, 0), 'BBCH Stage', font_size=9, bold=True,
               color=CLR_LABEL, align=PP_ALIGN.LEFT, fill_color=CLR_HEADER_BG)

    bbch_codes = crop_data.get('bbch_codes', [])
    for col_idx, stage_num in enumerate(existing_stages):
        code = bbch_codes[stage_num - 1] if (stage_num - 1) < len(bbch_codes) else ''
        style_cell(table.cell(0, col_idx + 1), code, font_size=10, bold=True,
                   color=CLR_BBCH, fill_color=CLR_HEADER_BG)

    # Row 1: Description
    style_cell(table.cell(1, 0), 'Description', font_size=9, bold=True,
               color=CLR_LABEL, align=PP_ALIGN.LEFT)

    descriptions = crop_data.get('descriptions', [])
    for col_idx, stage_num in enumerate(existing_stages):
        desc = descriptions[stage_num - 1] if (stage_num - 1) < len(descriptions) else ''
        style_cell(table.cell(1, col_idx + 1), desc, font_size=7, color=CLR_DESC)

    # Row 2: Your Product (empty with placeholder)
    style_cell(table.cell(2, 0), 'Your Product', font_size=9, bold=True,
               color=CLR_LABEL, align=PP_ALIGN.LEFT)

    for col_idx in range(num_stages):
        style_cell(table.cell(2, col_idx + 1), 'Add product\n& dosage',
                   font_size=6, color=CLR_PLACEHOLDER)

    # Set table borders
    from pptx.oxml.ns import qn
    tbl = table._tbl
    for row_idx in range(num_rows):
        for col_idx in range(num_cols):
            cell = table.cell(row_idx, col_idx)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                border = tcPr.find(qn(border_name))
                if border is None:
                    border = tcPr.makeelement(qn(border_name), {})
                    tcPr.append(border)

                border.set('w', str(Pt(0.5)))
                border.set('cmpd', 'sng')

                solidFill = border.find(qn('a:solidFill'))
                if solidFill is None:
                    solidFill = border.makeelement(qn('a:solidFill'), {})
                    border.append(solidFill)
                else:
                    solidFill.clear()

                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'C8CEC0'})
                solidFill.append(srgbClr)

    y_cursor += table_height + Inches(0.15)

    # ─── FOOTER ─────────────────────────────────────────────
    if include_footer:
        footer_box = slide.shapes.add_textbox(
            MARGIN_LEFT, SLIDE_HEIGHT - Inches(0.45),
            CONTENT_WIDTH, Inches(0.3)
        )
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'crop-stages.github.io'
        p.font.size = Pt(8)
        p.font.color.rgb = CLR_FOOTER
        p.font.name = 'Segoe UI'
        p.alignment = PP_ALIGN.RIGHT

    if images_found == 0:
        print(f"  ⚠ No images found for '{crop_slug}'")
    else:
        print(f"  ✓ {images_found}/{num_stages} images added")

    return slide


def process_single_html(html_path, images_dir, output_dir):
    """Process one HTML file → one PPTX file."""
    print(f"Processing: {os.path.basename(html_path)}")
    crop_data = parse_html(html_path)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_crop_slide(prs, crop_data, images_dir)

    stem = Path(html_path).stem
    output_path = os.path.join(output_dir, f'{stem}.pptx')
    prs.save(output_path)
    print(f"  → Saved: {output_path}")
    return output_path


def process_all_to_single(html_dir, images_dir, output_path):
    """Process all HTML files in a directory → one PPTX with multiple slides."""
    html_files = sorted(
        glob.glob(os.path.join(html_dir, '*.html'))
    )

    if not html_files:
        print(f"No HTML files found in {html_dir}")
        return

    print(f"Found {len(html_files)} HTML files")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for html_path in html_files:
        print(f"\nProcessing: {os.path.basename(html_path)}")
        try:
            crop_data = parse_html(html_path)
            add_crop_slide(prs, crop_data, images_dir)
        except Exception as e:
            print(f"  ERROR: {e}")

    prs.save(output_path)
    print(f"\n✅ Saved combined PPTX: {output_path} ({len(html_files)} slides)")


def main():
    parser = argparse.ArgumentParser(
        description='Convert HTML crop growth stage tables to PowerPoint'
    )
    parser.add_argument('--html-dir', required=True,
                        help='Directory with crop HTML files')
    parser.add_argument('--images-dir', required=True,
                        help='Directory with crop images (e.g. assets/images/crops)')
    parser.add_argument('--output-dir', default='./pptx_output',
                        help='Output directory for PPTX files')
    parser.add_argument('--single-file', action='store_true',
                        help='Combine all crops into one PPTX')
    parser.add_argument('--single-file-name', default='all_crops.pptx',
                        help='Name for combined PPTX file')

    args = parser.parse_args()

    os.makedirs(args.output_dir, exist_ok=True)

    if args.single_file:
        output_path = os.path.join(args.output_dir, args.single_file_name)
        process_all_to_single(args.html_dir, args.images_dir, output_path)
    else:
        html_files = sorted(glob.glob(os.path.join(args.html_dir, '*.html')))
        print(f"Found {len(html_files)} HTML files")
        for html_path in html_files:
            try:
                process_single_html(html_path, args.images_dir, args.output_dir)
            except Exception as e:
                print(f"  ERROR processing {html_path}: {e}")


if __name__ == '__main__':
    main()
